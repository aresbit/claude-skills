#!/usr/bin/env python3
"""Batch convert slides/PDFs to PDF + Markdown.

Supported input extensions:
- slides: .ppt, .pptx, .pps, .ppsx, .odp
- pdf: .pdf (text-based only; no OCR)
"""

from __future__ import annotations

import argparse
import shutil
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import List, Optional, Tuple


SLIDE_EXTS = {".ppt", ".pptx", ".pps", ".ppsx", ".odp"}
PDF_EXTS = {".pdf"}
SUPPORTED_EXTS = SLIDE_EXTS | PDF_EXTS


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Batch convert slides/PDFs to PDF + Markdown."
    )
    parser.add_argument("input_dir", type=Path, help="Directory containing slide files")
    parser.add_argument(
        "--output-dir",
        type=Path,
        default=Path("./slide-export"),
        help="Output root directory (default: ./slide-export)",
    )
    parser.add_argument(
        "--recursive",
        action="store_true",
        help="Recursively search subdirectories",
    )
    parser.add_argument(
        "--overwrite",
        action="store_true",
        help="Overwrite existing output files",
    )
    return parser.parse_args()


def ensure_soffice() -> str:
    soffice = shutil.which("soffice")
    if not soffice:
        raise RuntimeError("LibreOffice 'soffice' not found in PATH.")
    return soffice


def find_input_files(root: Path, recursive: bool) -> List[Path]:
    pattern = "**/*" if recursive else "*"
    files = []
    for path in root.glob(pattern):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTS:
            files.append(path)
    return sorted(files)


def run_soffice_convert(
    soffice: str,
    source_file: Path,
    target_filter: str,
    outdir: Path,
) -> Path:
    outdir.mkdir(parents=True, exist_ok=True)
    cmd = [
        soffice,
        "--headless",
        "--convert-to",
        target_filter,
        "--outdir",
        str(outdir),
        str(source_file),
    ]
    result = subprocess.run(cmd, capture_output=True, text=True, check=False)
    if result.returncode != 0:
        raise RuntimeError(
            f"soffice convert failed for {source_file}\n"
            f"stdout: {result.stdout}\n"
            f"stderr: {result.stderr}"
        )
    output_suffix = "." + target_filter.split(":")[0]
    candidate = outdir / f"{source_file.stem}{output_suffix}"
    if not candidate.exists():
        alt_candidates = sorted(outdir.glob(f"{source_file.stem}*{output_suffix}"))
        if not alt_candidates:
            raise RuntimeError(f"Cannot find converted file for {source_file}")
        candidate = alt_candidates[0]
    return candidate


def extract_slide_text(slide) -> List[str]:
    lines: List[str] = []
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        text = shape.text_frame.text.strip()
        if text:
            lines.append(text)
    return lines


def extract_notes(slide) -> str:
    notes = getattr(slide, "notes_slide", None)
    if not notes:
        return ""
    frame = getattr(notes.notes_text_frame, "text", "")
    return frame.strip() if frame else ""


def write_markdown_from_pptx(pptx_file: Path, md_file: Path, title_hint: str) -> None:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise RuntimeError(
            "Missing dependency: python-pptx. Install with: pip install python-pptx"
        ) from exc

    prs = Presentation(str(pptx_file))
    md_file.parent.mkdir(parents=True, exist_ok=True)

    parts: List[str] = [f"# {title_hint}", ""]
    for idx, slide in enumerate(prs.slides, start=1):
        parts.append(f"## Slide {idx}")
        parts.append("")

        text_blocks = extract_slide_text(slide)
        if text_blocks:
            for block in text_blocks:
                for line in block.splitlines():
                    cleaned = line.strip()
                    if cleaned:
                        parts.append(f"- {cleaned}")
        else:
            parts.append("- (No visible text)")

        notes = extract_notes(slide)
        if notes:
            parts.append("")
            parts.append("### Notes")
            parts.append("")
            for line in notes.splitlines():
                line = line.strip()
                if line:
                    parts.append(f"> {line}")
        parts.append("")

    md_file.write_text("\n".join(parts).rstrip() + "\n", encoding="utf-8")


def write_markdown_from_pdf(pdf_file: Path, md_file: Path, title_hint: str) -> None:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise RuntimeError("Missing dependency: pypdf. Install with: pip install pypdf") from exc

    reader = PdfReader(str(pdf_file))
    parts: List[str] = [f"# {title_hint}", ""]
    non_empty_pages = 0

    for idx, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        parts.append(f"## Page {idx}")
        parts.append("")
        if text:
            non_empty_pages += 1
            for para in text.splitlines():
                line = para.strip()
                if line:
                    parts.append(line)
        else:
            parts.append("- (No extractable text on this page)")
        parts.append("")

    if non_empty_pages == 0:
        parts.append("> This PDF appears image-based or has no embedded text. OCR is not included.")
        parts.append("")

    md_file.parent.mkdir(parents=True, exist_ok=True)
    md_file.write_text("\n".join(parts).rstrip() + "\n", encoding="utf-8")


def convert_one_file(
    src: Path,
    input_root: Path,
    output_root: Path,
    soffice: Optional[str],
    overwrite: bool,
) -> Tuple[Path, Path]:
    rel = src.relative_to(input_root)
    rel_no_ext = rel.with_suffix("")
    pdf_out = output_root / "pdf" / rel_no_ext.with_suffix(".pdf")
    md_out = output_root / "md" / rel_no_ext.with_suffix(".md")

    pdf_out.parent.mkdir(parents=True, exist_ok=True)
    md_out.parent.mkdir(parents=True, exist_ok=True)

    if (not overwrite) and pdf_out.exists() and md_out.exists():
        return pdf_out, md_out

    ext = src.suffix.lower()
    if ext in PDF_EXTS:
        if overwrite or not pdf_out.exists():
            shutil.copy2(src, pdf_out)
        write_markdown_from_pdf(src, md_out, src.stem)
        return pdf_out, md_out

    if ext in SLIDE_EXTS:
        if not soffice:
            raise RuntimeError("LibreOffice 'soffice' is required for slide conversion.")

        with tempfile.TemporaryDirectory(prefix="slide-convert-") as tmp:
            tmpdir = Path(tmp)
            converted_pdf = run_soffice_convert(soffice, src, "pdf", tmpdir)
            if overwrite or not pdf_out.exists():
                shutil.copy2(converted_pdf, pdf_out)

            pptx_for_md: Optional[Path] = None
            if ext == ".pptx":
                pptx_for_md = src
            else:
                try:
                    converted_pptx = run_soffice_convert(
                        soffice, src, "pptx:Impress MS PowerPoint 2007 XML", tmpdir
                    )
                    pptx_for_md = converted_pptx
                except Exception:
                    pptx_for_md = None

            if pptx_for_md is None:
                md_out.write_text(
                    f"# {src.stem}\n\n- Markdown extraction skipped: unable to produce PPTX intermediate.\n",
                    encoding="utf-8",
                )
            else:
                write_markdown_from_pptx(pptx_for_md, md_out, src.stem)
        return pdf_out, md_out

    raise RuntimeError(f"Unsupported extension: {ext}")


def main() -> int:
    args = parse_args()

    input_dir = args.input_dir.resolve()
    output_dir = args.output_dir.resolve()

    if not input_dir.exists() or not input_dir.is_dir():
        print(f"Input directory not found: {input_dir}", file=sys.stderr)
        return 2

    files = find_input_files(input_dir, args.recursive)
    if not files:
        print(f"No supported files found in {input_dir}", file=sys.stderr)
        return 1

    needs_soffice = any(f.suffix.lower() in SLIDE_EXTS for f in files)
    soffice: Optional[str] = None
    if needs_soffice:
        try:
            soffice = ensure_soffice()
        except RuntimeError as exc:
            print(str(exc), file=sys.stderr)
            return 2

    ok_count = 0
    failed: List[Tuple[Path, str]] = []

    for src in files:
        try:
            pdf_out, md_out = convert_one_file(
                src=src,
                input_root=input_dir,
                output_root=output_dir,
                soffice=soffice,
                overwrite=args.overwrite,
            )
            ok_count += 1
            print(f"[ok] {src} -> {pdf_out} | {md_out}")
        except Exception as exc:
            failed.append((src, str(exc)))
            print(f"[failed] {src}: {exc}", file=sys.stderr)

    print(f"\nDone. success={ok_count} failed={len(failed)}")
    if failed:
        print("Failed files:", file=sys.stderr)
        for src, err in failed:
            print(f"- {src}: {err}", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
